"""Parse uploaded files into a list of (text, page_number) segments.

Per-format behavior:
  - txt / cs / md : one segment, page_number=None
  - pdf           : one segment per page, page_number=<1-based>
  - pptx          : one segment per slide, page_number=<slide index>
  - docx          : one segment, page_number=None (no native page concept)
"""
from __future__ import annotations

import io
from dataclasses import dataclass

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class ParsedSegment:
    text: str
    page_number: int | None  # 1-based; None for non-paginated formats


def _parse_pdf(data: bytes) -> list[ParsedSegment]:
    reader = PdfReader(io.BytesIO(data))
    out: list[ParsedSegment] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            out.append(ParsedSegment(text=text, page_number=i))
    return out


def _parse_docx(data: bytes) -> list[ParsedSegment]:
    """Flatten paragraphs + table cells into a single segment.

    Word documents have no exposed page concept (page breaks are layout
    decisions made at render time), so we emit one segment for the whole
    file. Tables are joined with " | " between cells so chunks retain
    some readable structure.
    """
    doc = DocxDocument(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return [ParsedSegment(text=text, page_number=None)] if text.strip() else []


def _parse_pptx(data: bytes) -> list[ParsedSegment]:
    """One segment per slide, page_number = slide index.

    Walks every shape on the slide that has a text frame; concatenates
    paragraph text with newlines. Speaker notes are intentionally excluded
    — they're authoring metadata, not deck content.
    """
    prs = Presentation(io.BytesIO(data))
    out: list[ParsedSegment] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    parts.append(t)
        text = "\n".join(parts)
        if text.strip():
            out.append(ParsedSegment(text=text, page_number=i))
    return out


def parse(extension: str, data: bytes) -> list[ParsedSegment]:
    if extension in ("txt", "cs", "md"):
        return [ParsedSegment(text=data.decode("utf-8", errors="replace"), page_number=None)]
    if extension == "pdf":
        return _parse_pdf(data)
    if extension == "docx":
        return _parse_docx(data)
    if extension == "pptx":
        return _parse_pptx(data)
    raise ValueError(f"unsupported extension: {extension}")
